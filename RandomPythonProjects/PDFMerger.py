from pptx import Presentation

# Load the first two presentations
prs1 = Presentation("1.pptx")
prs2 = Presentation("2.pptx")

# Create a new presentation to hold the combined slides
prs3 = Presentation()

# Loop through each slide in the first two presentations
for slide in prs1.slides + prs2.slides:
    # Copy the slide to the new presentation
    new_slide = prs3.slides.add_slide(slide.slide_layout)
    for shape in slide.shapes:
        new_shape = new_slide.shapes.element.add_shape(shape.shape_type.__name__, shape.left, shape.top, shape.width, shape.height)
        new_shape.text_frame.text = shape.text

# Open the third presentation and add its slides to the new presentation
prs4 = Presentation("3.pptx")
for slide in prs4.slides:
    prs3.slides.add_slide(slide.slide_layout)

# Save the combined presentation
prs3.save("combined_presentation.pptx")